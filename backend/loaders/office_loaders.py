"""Text extraction for PDF, DOCX, PPTX, CSV, and XLSX files."""

from __future__ import annotations

import io

import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from PyPDF2 import PdfReader


def parse_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def parse_docx(data: bytes) -> str:
    doc = DocxDocument(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                parts.append(row_text)
    return "\n".join(parts)


def parse_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    slides_text = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines = [f"## Slide {idx}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        lines.append(row_text)
        if len(lines) > 1:
            slides_text.append("\n".join(lines))
    return "\n\n".join(slides_text)


def _dataframe_to_text(df: pd.DataFrame) -> str:
    """Render a dataframe as 'column: value' rows so embeddings capture column semantics."""
    df = df.fillna("")
    return "\n".join("; ".join(f"{col}: {row[col]}" for col in df.columns) for _, row in df.iterrows())


def parse_csv(data: bytes) -> str:
    df = pd.read_csv(io.BytesIO(data))
    return _dataframe_to_text(df)


def parse_xlsx(data: bytes) -> str:
    sheets = pd.read_excel(io.BytesIO(data), sheet_name=None)
    return "\n\n".join(f"## Sheet: {name}\n{_dataframe_to_text(df)}" for name, df in sheets.items())
